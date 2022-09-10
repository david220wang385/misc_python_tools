"""
References:
https://stackoverflow.com/questions/39418620/extracting-text-from-multiple-powerpoint-files-using-python
https://stackoverflow.com/questions/29110950/python-concordance-command-in-nltk
https://stackoverflow.com/questions/6416131/add-a-new-item-to-a-dictionary-in-python
https://stackoverflow.com/questions/29778519/stop-concordance-printing-no-matches-in-python
https://www.tutorialspoint.com/ternary-operator-in-python
https://www.nltk.org/api/nltk.html?highlight=concordance
https://stackoverflow.com/questions/29613487/multiple-lines-in-python-argparse-help-display
"""

import os, argparse, glob
from pptx import Presentation
import nltk.corpus  
from nltk.text import Text  

# Entrypoint
def main():

    # Provide decription of the program
    parser = argparse.ArgumentParser(
        description="""Search all PowerPoint files in a provided folder for text.
    \nThis program was cooked up by me on a whim in like an hour or so so its not really my best work
    -David Wang""",
        formatter_class=argparse.RawTextHelpFormatter
    )

    # Positional arguments for specifying target directory/file and key to be used in the keystream
    parser.add_argument(
        '--dir', 
        type=str, 
        default='',
        help='The directory with PPTs that you want to search, make sure to use "./" if you are selecting a folder in the current working directory'

    )
    parser.add_argument(
        '--width', 
        type=int, 
        default=80,
        help='Amount of context when the text matches the searched pattern'
    )

    # Parse arguments and initialize variables
    args = parser.parse_args()
    dir = args.dir
    width = args.width

    # Printing for debug
    print(" ARGUMENTS PARSED: ".center(100, '-'))
    print("Selected directory:\t\t", "Current executing directory" if dir == "" else dir) 
    print("Concordance width:\t\t", width)

    corpus = {}  # Store all text across all powerpoints

    # Build a text_list of all ppt text before searching
    # Iterate thru each ppt file
    for eachfile in glob.glob(dir+"/*.pptx"):
        
        # Track all text in a single ppt
        ppt_text = ""
        ppt_text += eachfile

        prs = Presentation(eachfile)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    ppt_text += shape.text + " "

        corpus[eachfile] = ppt_text

    if len(corpus.items()) == 0:
        print("No .pptx files detected in the provided directory")
        exit()
    print("PPTs successfully loaded")
    print("".center(100, '-'))

    # Prompt user 
    user_input = " "
    while user_input != "":

        # Prompt user to enter string to search for in the .pptx files
        print()
        user_input = input("> ENTER TEXT TO SEARCH FOR: ")

        for (key, value) in corpus.items():
            text_list = nltk.ConcordanceIndex(nltk.word_tokenize(value))

            if text_list.offsets(user_input):
                print()
                print((" " + key[:80] + " ").center(100, '-'), "\n")
                text_list.print_concordance(user_input, width=width)
                print("".center(100, '-'))
        
        print()
        print("".center(100, '*'))


# Only execute if called directly
if __name__ == "__main__":
    main()